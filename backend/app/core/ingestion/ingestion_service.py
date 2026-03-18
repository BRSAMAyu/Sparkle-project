from __future__ import annotations
import base64
import io
import os
import re
from typing import Any

from fastapi import HTTPException
from loguru import logger
from pydantic import BaseModel

from app.config import settings
from app.services.ocr_service import ocr_service

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    Document = None
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    Presentation = None
    HAS_PPTX = False

try:
    from PIL import ImageOps
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

class ExtractedChunk(BaseModel):
    text: str
    page_num: int
    source: str
    metadata: dict = {}  # bold, header, color, etc.
    ocr_confidence: float | None = None  # 0.0-1.0, None if not from OCR

class IngestionService:
    """
    Robust Document Ingestion Service for Exam Savior.
    Handles PDF, DOCX, PPTX with advanced cleaning and metadata extraction.
    """

    def process_file(self, file_path: str, options: dict[str, Any] = None) -> list[ExtractedChunk]:
        """
        Main entry point. Dispatches to specific handlers based on extension.
        Includes Magic Byte validation.
        """
        if options is None:
            options = {}

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        self._validate_magic_bytes(file_path)

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        try:
            if ext == ".pdf":
                return self._process_pdf(file_path, options)
            elif ext == ".docx":
                return self._process_docx(file_path)
            elif ext == ".pptx":
                return self._process_pptx(file_path)
            elif ext in [".jpg", ".jpeg", ".png", ".webp", ".gif"]:
                return self._process_image(file_path, options)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return []
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            raise e

    def _validate_magic_bytes(self, file_path: str):
        """
        Check if the file header matches its extension to prevent spoofing.
        Enhanced with comprehensive magic byte validation.
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        with open(file_path, "rb") as f:
            # 读取前512字节用于魔数检测
            header = f.read(512)
            if not header:
                raise ValueError("Empty file or cannot read file header")

        # 验证魔数
        if ext == ".pdf":
            # PDF文件以 %PDF- 开头
            if not header.startswith(b"%PDF-"):
                raise ValueError(
                    "Invalid PDF file: missing PDF magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        elif ext in [".docx", ".xlsx", ".pptx"]:
            # Office文档都是ZIP格式，以PK\x03\x04开头
            if not header.startswith(b"PK\x03\x04"):
                raise ValueError(
                    "Invalid Office document: missing ZIP magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        elif ext == ".png":
            # PNG以 \x89PNG\r\n\x1a\n 开头
            png_magic = b'\x89PNG\r\n\x1a\n'
            if not header.startswith(png_magic):
                raise ValueError(
                    "Invalid PNG file: missing PNG magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        elif ext in [".jpg", ".jpeg"]:
            # JPEG以 \xFF\xD8\xFF 开头
            if not header.startswith(b'\xFF\xD8\xFF'):
                raise ValueError(
                    "Invalid JPEG file: missing JPEG magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        elif ext == ".gif":
            # GIF以 GIF87a 或 GIF89a 开头
            if not (header.startswith(b"GIF87a") or header.startswith(b"GIF89a")):
                raise ValueError(
                    "Invalid GIF file: missing GIF magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        elif ext == ".webp":
            # WebP以 RIFF....WEBP 开头
            if len(header) < 12:
                raise ValueError("File too short to be a valid WebP")
            if not header.startswith(b"RIFF") or header[8:12] != b"WEBP":
                raise ValueError(
                    "Invalid WebP file: missing WebP magic bytes. "
                    "File may be corrupted or renamed with wrong extension."
                )

        else:
            logger.warning(f"No magic byte validation implemented for extension: {ext}")

    def _process_pdf(self, path: str, options: dict[str, Any]) -> list[ExtractedChunk]:
        chunks = []
        enable_ocr = bool(options.get("enable_ocr", True))
        # Use pdfplumber for better layout analysis
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                ocr_confidence = None

                # --- OCR Fallback Strategy ---
                # If text is empty or suspiciously short (scanned page), try OCR
                if enable_ocr and len(text.strip()) < 50:
                    logger.info(f"Page {i+1} has low text content ({len(text.strip())} chars). Attempting OCR...")
                    ocr_text, ocr_confidence = self._attempt_ocr(page, options)
                    if ocr_text:
                        text = ocr_text
                        logger.info(
                            f"OCR recovered {len(text)} chars from Page {i+1} "
                            f"(confidence: {ocr_confidence:.2f})" if ocr_confidence is not None else
                            f"OCR recovered {len(text)} chars from Page {i+1}"
                        )
                    else:
                        logger.warning(f"OCR failed or produced no text for Page {i+1}")

                if not text:
                    continue

                # Cleaning
                clean_text = self._clean_text(text)

                if len(clean_text) < 20:  # Skip empty/noise pages
                    continue

                chunks.append(ExtractedChunk(
                    text=clean_text,
                    page_num=i + 1,
                    source="pdf",
                    metadata={"raw_len": len(text)},
                    ocr_confidence=ocr_confidence
                ))
        return chunks

    def _attempt_ocr(self, page, options: dict[str, Any]) -> tuple[str, float | None]:
        """
        Helper to run OCR on a pdfplumber page object.
        Dispatches to local Tesseract or Remote API based on options.
        """
        if not HAS_PIL:
            logger.warning("OCR requested but Pillow not installed.")
            return "", None

        try:
            # pdfplumber to_image returns a PageImage, .original gives PIL Image
            # Use 300 DPI for better OCR
            im = page.to_image(resolution=300).original

            ocr_engine = (options.get("ocr_engine", "local") or "local").lower()
            if ocr_engine == "deepseek":
                logger.warning("OCR engine 'deepseek' 已废弃，自动切换到 'zhipu'")
                ocr_engine = "zhipu"

            if ocr_engine == "zhipu":
                return self._ocr_via_api(im, options), None
            if not HAS_TESSERACT:
                logger.warning("Local OCR requested but pytesseract not installed.")
                return "", None
            return self._ocr_via_local(im)

        except Exception as e:
            logger.warning(f"OCR Error: {e}")
            return "", None

    def _ocr_via_local(self, im) -> tuple[str, float | None]:
        """Run local Tesseract OCR with preprocessing"""
        try:
            # --- Image Preprocessing for Accuracy ---
            # 1. Convert to grayscale
            im = im.convert('L')

            # 2. Auto-contrast (stretch histogram)
            im = ImageOps.autocontrast(im)

            # 3. Simple Binarization (Thresholding)
            threshold = 200
            im = im.point(lambda p: 255 if p > threshold else 0)

            # 4. Run OCR
            try:
                config = r'--oem 3 --psm 6'
                # Get detailed data including confidence
                ocr_data = pytesseract.image_to_data(
                    im, lang='chi_sim+eng', config=config, output_type=pytesseract.Output.DICT
                )
                text = pytesseract.image_to_string(im, lang='chi_sim+eng', config=config).strip()
                confidences = [
                    int(conf) for conf in ocr_data.get("conf", []) if conf.isdigit() and int(conf) >= 0
                ]
                avg_confidence = (sum(confidences) / len(confidences) / 100.0) if confidences else None
                return text, avg_confidence
            except pytesseract.TesseractError:
                config = r'--oem 3 --psm 6'
                return pytesseract.image_to_string(im, lang='eng', config=config).strip(), None
        except Exception as e:
            logger.warning(f"Local OCR Failed: {e}")
            return "", None

    def _ocr_via_api(self, image, options: dict[str, Any]) -> str:
        """
        Run remote GLM OCR.
        """
        del options
        if not settings.ZHIPU_API_KEY:
            logger.warning("GLM OCR requested but ZHIPU_API_KEY not set.")
            return ""

        try:
            # Convert to base64
            buffered = io.BytesIO()
            # Convert to RGB to ensure compatibility
            if image.mode != "RGB":
                image = image.convert("RGB")

            image.save(buffered, format="JPEG", quality=95)
            img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
            return ocr_service.ocr_from_base64_sync(img_str)

        except Exception as e:
            logger.error(f"GLM OCR API Exception: {e}")
            return ""

    def _process_docx(self, path: str) -> list[ExtractedChunk]:
        if not HAS_DOCX:
            raise HTTPException(
                status_code=501,
                detail="DOCX processing requires python-docx, which is not installed."
            )
        doc = Document(path)
        chunks = []
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue

            # Feature Engineering: Extract styles
            style_name = para.style.name.lower()
            is_header = "heading" in style_name

            # Check for bold/color runs
            is_bold = any(run.bold for run in para.runs)

            metadata = {
                "is_header": is_header,
                "is_bold": is_bold,
                "style": style_name
            }

            clean_text = self._clean_text(text)

            chunks.append(ExtractedChunk(
                text=clean_text,
                page_num=i, # Docx doesn't have strict pages, use para index
                source="docx",
                metadata=metadata
            ))
        return chunks

    def _process_pptx(self, path: str) -> list[ExtractedChunk]:
        if not HAS_PPTX:
            raise HTTPException(
                status_code=501,
                detail="PPTX processing requires python-pptx, which is not installed."
            )
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            # CRITICAL: Extract speaker notes (where the real content often lives)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    slide_text.append(f"[NOTES]: {notes}")

            full_text = "\n".join(slide_text)
            clean_text = self._clean_text(full_text)

            if len(clean_text) < 10:
                continue

            chunks.append(ExtractedChunk(
                text=clean_text,
                page_num=i + 1,
                source="pptx",
                metadata={}
            ))
        return chunks

    def _process_image(self, path: str, options: dict[str, Any] = None) -> list[ExtractedChunk]:
        """
        处理图片文件，直接进行 OCR。
        支持格式: JPG, JPEG, PNG, WebP, GIF
        """
        if options is None:
            options = {}

        if not HAS_PIL:
            raise HTTPException(
                status_code=501,
                detail="Image OCR requires Pillow, which is not installed."
            )

        from PIL import Image

        try:
            im = Image.open(path)

            # 获取图片格式信息
            image_format = im.format or "UNKNOWN"
            logger.info(f"Processing image: {path}, format={image_format}, size={im.size}")

            # 创建 Mock Page 对象用于复用现有 OCR 逻辑
            class MockPage:
                def __init__(self, image):
                    self._image = image

                def to_image(self, resolution=300):
                    class PageImage:
                        def __init__(self, img):
                            self.original = img
                    return PageImage(self._image)

            mock_page = MockPage(im)

            # 复用现有 OCR 逻辑
            enable_ocr = bool(options.get("enable_ocr", True))
            if not enable_ocr:
                logger.warning("Image processing requires OCR, but OCR is disabled")
                return []

            ocr_text, ocr_confidence = self._attempt_ocr(mock_page, options)

            if not ocr_text:
                logger.warning(f"OCR produced no text for image: {path}")
                return []

            # 清理文本
            clean_text = self._clean_text(ocr_text)

            if len(clean_text) < 20:
                logger.info(f"Image OCR text too short ({len(clean_text)} chars), skipping")
                return []

            return [ExtractedChunk(
                text=clean_text,
                page_num=1,
                source="image",
                metadata={
                    "format": image_format,
                    "width": im.size[0],
                    "height": im.size[1],
                },
                ocr_confidence=ocr_confidence
            )]

        except Exception as e:
            logger.error(f"Failed to process image {path}: {e}")
            raise

    def _clean_text(self, text: str) -> str:
        """
        Applies cleaning rules defined in Protocol.
        """
        # 1. De-hyphenation (simple heuristic)
        # Fix "exam-\nple" -> "example"
        text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)

        # 2. Remove Page Numbers (standalone numbers on lines)
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)

        # 3. Collapse multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)

        # 4. Remove common watermarks (example)
        text = text.replace("Do Not Distribute", "")

        # 5. Fix common encoding artifacts
        text = text.replace("\x00", "") # Null bytes

        return text.strip()

# Singleton instance
ingestion_service = IngestionService()
