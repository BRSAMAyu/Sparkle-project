#!/usr/bin/env python3
"""
Sparkle V2 Roadshow PPT Generator
11 main slides (8-min presentation) + 8 appendix slides (Q&A defense)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy for headings
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)   # Body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF7, 0xF8, 0xFA)   # Page background
BLUE      = RGBColor(0x25, 0x63, 0xEB)    # Primary accent
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)    # Darker blue
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)    # Warm accent
GRAY      = RGBColor(0x6B, 0x72, 0x80)    # Secondary text
LIGHT_GRAY = RGBColor(0xE2, 0xE5, 0xEA)   # Borders
GREEN     = RGBColor(0x10, 0xB9, 0x81)    # Success
RED       = RGBColor(0xEF, 0x44, 0x44)    # Danger/failure
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)    # Card white
SOFT_BLUE = RGBColor(0xDB, 0xE4, 0xFE)    # Light blue bg
SOFT_ORANGE = RGBColor(0xFE, 0xF3, 0xC7)  # Light orange bg
SOFT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)   # Light green bg

# ── Slide Dimensions (16:9 widescreen) ─────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
BLANK_LAYOUT = prs.slide_layouts[6]  # blank

# ── Helper Functions ───────────────────────────────────────────

def add_blank_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="",
                font_size=18, font_color=DARK_TEXT, bold=False,
                alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei',
                line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return tf

def add_rich_textbox(slide, left, top, width, height):
    """Returns text_frame for multi-paragraph content"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=16, font_color=DARK_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei',
             space_after=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p

def add_shape(slide, left, top, width, height, fill_color=None,
              border_color=None, border_width=None, shape_type=MSO_SHAPE.RECTANGLE,
              corner_radius=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     border_color=None, border_width=None):
    shape = add_shape(slide, left, top, width, height, fill_color,
                      border_color, border_width, MSO_SHAPE.ROUNDED_RECTANGLE)
    return shape

def add_circle(slide, left, top, size, fill_color=None, border_color=None):
    shape = add_shape(slide, left, top, size, size, fill_color,
                      border_color, None, MSO_SHAPE.OVAL)
    return shape

def add_page_number(slide, num, total=11):
    add_textbox(slide, 12.2, 7.05, 1.0, 0.35, f"{num} / {total}",
                font_size=9, font_color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def add_section_label(slide, text, top=0.3):
    """Small section label at top of page"""
    add_textbox(slide, 0.6, top, 3.0, 0.35, text, font_size=10,
                font_color=BLUE, bold=True)

def add_title_bar(slide, title_text, subtitle_text=None, page_num=None, total=11):
    """Standard page title"""
    # Top accent line
    add_shape(slide, 0, 0, 13.333, 0.06, BLUE)
    # Title
    add_textbox(slide, 0.8, 0.3, 11.5, 0.7, title_text,
                font_size=32, font_color=NAVY, bold=True)
    if subtitle_text:
        add_textbox(slide, 0.8, 0.9, 11.5, 0.4, subtitle_text,
                    font_size=14, font_color=GRAY)
    if page_num:
        add_page_number(slide, page_num, total)

def add_phone_frame(slide, left, top, width=2.6, height=5.2):
    """Draw a phone mockup frame"""
    phone = add_rounded_rect(slide, left, top, width, height,
                              WHITE, LIGHT_GRAY, 2)
    # Notch area
    add_shape(slide, left + width/2 - 0.3, top + 0.1, 0.6, 0.15,
              DARK_TEXT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    return phone

def add_task_card_mockup(slide, left, top, width, height, task_data):
    """Draw a task card UI element with shapes"""
    card = add_rounded_rect(slide, left, top, width, height, WHITE, LIGHT_GRAY, 1.5)

    # Status bar
    add_shape(slide, left + 0.15, top + 0.15, width - 0.3, 0.06, BLUE)

    # Task title
    add_textbox(slide, left + 0.2, top + 0.3, width - 0.4, 0.35,
                task_data.get('title', 'Task Card'),
                font_size=13, font_color=NAVY, bold=True)

    # Sections
    y = top + 0.7
    for label, value in [('Why', task_data.get('why', '')),
                          ('Materials', task_data.get('materials', '')),
                          ('Steps', task_data.get('steps', '')),
                          ('Stuck?', task_data.get('stuck', ''))]:
        if value:
            add_textbox(slide, left + 0.2, y, 0.6, 0.2, label,
                        font_size=8, font_color=BLUE, bold=True)
            add_textbox(slide, left + 0.8, y, width - 1.1, 0.2, value,
                        font_size=8, font_color=DARK_TEXT)
            y += 0.22

    # Progress indicator
    add_shape(slide, left + 0.2, top + height - 0.35, width - 0.4, 0.04, LIGHT_GRAY)
    add_shape(slide, left + 0.2, top + height - 0.35, (width - 0.4) * 0.6, 0.04, BLUE)

    return card


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: HOOK — 7 Days to Exam
# ═══════════════════════════════════════════════════════════════
def build_slide_1():
    slide = add_blank_slide()
    set_slide_bg(slide, NAVY)

    # Dramatic countdown
    add_textbox(slide, 0.8, 1.2, 5.0, 1.0, "7",
                font_size=120, font_color=WHITE, bold=True)
    add_textbox(slide, 2.8, 1.5, 5.0, 0.6, "天后考试",
                font_size=36, font_color=WHITE)
    add_textbox(slide, 0.8, 2.4, 8.0, 0.5, "基本没学。资料一堆。",
                font_size=24, font_color=RGBColor(0xBB, 0xBB, 0xCC))

    # The question
    add_textbox(slide, 0.8, 4.5, 10.0, 0.8, "AI 能救吗？",
                font_size=52, font_color=ORANGE, bold=True)

    # Bottom line
    add_textbox(slide, 0.8, 5.8, 10.0, 0.5,
                "ChatGPT 能答题，但没法保证他把复习走完。",
                font_size=18, font_color=GRAY)

    # Sparkle logo placeholder bottom-right
    add_textbox(slide, 9.5, 6.5, 3.5, 0.5, "Sparkle",
                font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, 9.5, 6.95, 3.5, 0.35, "鸿雁杯 · 2026",
                font_size=12, font_color=GRAY, alignment=PP_ALIGN.RIGHT)

    add_speaker_notes(slide, """[25 秒]

开场不要介绍项目名，直接把评委拉进场景。

口播：
"一个学生，7天后计网考试。基本没学。教材、真题、笔记堆了一桌。他打开 ChatGPT 问：我要考计网了，怎么办。ChatGPT 能给他一份复习计划。能解释知识点。能答疑。但它没法保证他把复习真正走完。这就是我们的起点。"

关键动作：第一秒就制造画面感。语气平静但紧迫。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM — Four Fractures
# ═══════════════════════════════════════════════════════════════
def build_slide_2():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "AI 很强，但普通人离目标没有更近", page_num=2)

    # Lead sentence
    add_textbox(slide, 0.8, 1.3, 11.5, 0.5,
                "面对一个真实目标，缺的不是答案，而是四件 AI 做不到的事：",
                font_size=18, font_color=DARK_TEXT)

    # Four fracture cards
    fractures = [
        ("目标 → 路径", "不知道从哪里开始\n什么该做、什么该放弃", BLUE, SOFT_BLUE),
        ("路径 → 执行", "做了但不知道对不对\n卡住了不知道问谁", ORANGE, SOFT_ORANGE),
        ("执行 → 反馈", "任务失败了就结束了\n不知道为什么失败", RED, RGBColor(0xFE, 0xE2, 0xE2)),
        ("反馈 → 修正", "下次同样的问题\n还是同样卡住", RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xED, 0xE9, 0xFE)),
    ]

    for i, (title, desc, accent, bg) in enumerate(fractures):
        x = 0.6 + i * 3.1
        y = 2.2
        # Card bg
        add_rounded_rect(slide, x, y, 2.8, 2.8, bg)
        # Accent top bar
        add_shape(slide, x, y, 2.8, 0.06, accent)
        # Number
        add_textbox(slide, x + 0.25, y + 0.35, 0.5, 0.4, str(i+1),
                    font_size=28, font_color=accent, bold=True)
        # Title
        add_textbox(slide, x + 0.25, y + 0.8, 2.3, 0.5, title,
                    font_size=16, font_color=NAVY, bold=True)
        # Desc
        add_textbox(slide, x + 0.25, y + 1.3, 2.3, 1.2, desc,
                    font_size=12, font_color=DARK_TEXT)

        # Arrow between cards (except last)
        if i < 3:
            add_textbox(slide, x + 2.85, y + 1.1, 0.4, 0.4, "→",
                        font_size=20, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Strategic choice - bottom
    add_shape(slide, 0.6, 5.5, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 5.8, 11.5, 0.5,
                "所有人都在卷模型能力。我们选择卷另一件事：让普通人真的把目标做成。",
                font_size=18, font_color=BLUE, bold=True)

    add_textbox(slide, 0.8, 6.4, 11.5, 0.4,
                "通用 AI 解决\"会不会答\"。Sparkle 解决\"能不能做成\"。",
                font_size=14, font_color=GRAY)

    add_speaker_notes(slide, """[40 秒]

口播：
"AI 这么强，为什么一个学生还是没法靠它通过考试？因为面对一个真实目标，缺的不是答案，而是四个东西：
第一，不知道怎么把目标变成可执行的路径。
第二，执行的时候不知道自己做得对不对。
第三，失败了没有反馈——任务打叉就结束了，不知道为什么失败。
第四，下次遇到同样的问题，还是卡住。

所有人都在卷模型能力。我们选择卷另一件事——让普通人真的把目标做成。

一句话：通用 AI 解决'会不会答'，Sparkle 解决'能不能做成'。"

过渡：下一页讲 Sparkle 第一次对话就体现出不同。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: AURORA — Decision, Not Answer
# ═══════════════════════════════════════════════════════════════
def build_slide_3():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "第一次对话就不一样：不是先回答，是先决策",
                  "Aurora 在每次对话前，先理解目标、状态、资料和偏好，再决定怎么回应",
                  page_num=3)

    # Left: Normal AI
    left_x = 0.6
    card_w = 5.5
    # Card header
    add_shape(slide, left_x, 1.5, card_w, 0.5, GRAY)
    add_textbox(slide, left_x + 0.2, 1.55, card_w - 0.4, 0.4,
                "普通 AI：收到问题 → 直接回答", font_size=14, font_color=WHITE, bold=True)

    # Chat bubble - user
    add_rounded_rect(slide, left_x + 3.2, 2.3, 2.0, 0.6, SOFT_BLUE)
    add_textbox(slide, left_x + 3.35, 2.35, 1.7, 0.5,
                "\"7天后考计网\n基本没学\"", font_size=10, font_color=DARK_TEXT, alignment=PP_ALIGN.LEFT)

    # Chat bubble - AI response (bad)
    add_rounded_rect(slide, left_x + 0.2, 2.3, 2.8, 2.5, RGBColor(0xF0, 0xF0, 0xF5))
    resp_lines = [
        "→ 推荐一份完整复习计划",
        "→ 包含全部 8 章内容",
        "→ 解释重点知识点",
        "→ 推荐参考书和视频",
        "",
        "❌ 不知道学生基础",
        "❌ 不知道哪些能放弃",
        "❌ 不知道 7 天够不够",
    ]
    tf = add_rich_textbox(slide, left_x + 0.35, 2.45, 2.5, 2.2)
    for i, line in enumerate(resp_lines):
        add_para(tf, line, font_size=9,
                 font_color=RED if line.startswith("❌") else GRAY,
                 bold=line.startswith("❌"), space_after=2)

    # VS divider
    add_textbox(slide, 6.4, 2.8, 0.5, 0.5, "VS",
                font_size=28, font_color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Right: Sparkle
    right_x = 7.2
    add_shape(slide, right_x, 1.5, card_w, 0.5, BLUE)
    add_textbox(slide, right_x + 0.2, 1.55, card_w - 0.4, 0.4,
                "Sparkle：先判断场景 → 再决定回应方式",
                font_size=14, font_color=WHITE, bold=True)

    # Chat bubble - same user input
    add_rounded_rect(slide, right_x + 3.2, 2.3, 2.0, 0.6, SOFT_BLUE)
    add_textbox(slide, right_x + 3.35, 2.35, 1.7, 0.5,
                "\"7天后考计网\n基本没学\"", font_size=10, font_color=DARK_TEXT)

    # Chat bubble - Sparkle response (good)
    add_rounded_rect(slide, right_x + 0.2, 2.3, 2.8, 2.5, SOFT_GREEN)
    resp_lines = [
        "→ 识别：7天抢救模式",
        "→ 先不推荐资料",
        "→ 要求做 12 分钟诊断",
        "→ 判断基础+考试权重",
        "",
        "✓ 知道该先测哪三块",
        "✓ 知道该放弃哪部分",
        "✓ 每一步基于决策",
    ]
    tf = add_rich_textbox(slide, right_x + 0.35, 2.45, 2.5, 2.2)
    for i, line in enumerate(resp_lines):
        add_para(tf, line, font_size=9,
                 font_color=GREEN if line.startswith("✓") else DARK_TEXT,
                 bold=line.startswith("✓"), space_after=2)

    # Bottom insight
    add_shape(slide, 0.6, 5.5, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 5.8, 11.5, 0.5,
                "Aurora 不是人格，是自适应决策层：不裸调模型，先组装上下文再回应。",
                font_size=16, font_color=NAVY, bold=True)
    add_textbox(slide, 0.8, 6.35, 11.5, 0.4,
                "别人把 AI 能力交给用户自己调。Sparkle 用 Aurora 帮用户自动调好。",
                font_size=14, font_color=GRAY)

    add_speaker_notes(slide, """[40 秒]

这是全 deck 最重要的"认知转换"页。不要讲"Aurora 回答更好"，要证明"决策入口不同"。

口播：
"当学生说'7 天后考试，基本没学'——
普通 AI 会怎么做？他给一份完整复习计划。8 章全覆盖。推荐资料。解释知识点。
但这份计划不知道学生基础为零。不知道 7 天不够学完 8 章。不知道哪些内容可以放弃。

Sparkle 怎么做？它先不回答。它先做判断：这是抢救模式。基础未知。需要低成本第一步。
所以它不做诊断前不给计划。它要先测三块投入产出比最高的内容。

这不是先问两个问题再回答。这是决策系统在运行。
我们给这套决策系统起了个名字叫 Aurora。它让 Sparkle 每次对话前先理解目标、状态、资料和偏好，再决定怎么回应。别人把 AI 能力交给用户自己调，我们帮用户自动调好。"

过渡：下面用一个真实案例证明这不是理论。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: RESCUE PART 1 — Diagnosis & Triage
# ═══════════════════════════════════════════════════════════════
def build_slide_4():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "7 天救急 · 上：诊断 + 主动放弃",
                  "极限压力测试：时间少、基础差、资料乱、必须取舍",
                  page_num=4)

    # Timeline visual
    # Day 0 marker
    add_circle(slide, 0.8, 2.0, 0.5, BLUE)
    add_textbox(slide, 0.85, 2.05, 0.4, 0.4, "0", font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.45, 2.05, 1.5, 0.4, "第 0 天", font_size=14, font_color=BLUE, bold=True)

    # Timeline line
    add_shape(slide, 1.05, 2.5, 0.03, 3.2, BLUE)

    # Moment 1: Diagnosis
    add_circle(slide, 0.9, 2.8, 0.3, BLUE)
    add_rounded_rect(slide, 1.55, 2.65, 5.2, 2.2, CARD_BG, LIGHT_GRAY, 1.5)
    add_textbox(slide, 1.75, 2.75, 4.8, 0.3, "时刻 1：12 分钟诊断",
                font_size=15, font_color=NAVY, bold=True)
    diag_text = [
        "✓ 第 1-3 章：掌握度 75% — 快速过即可",
        "✓ 第 4-6 章：掌握度 30% — 需要重点投入",
        "✓ 第 7-8 章：掌握度 5% — 几乎零基础",
        "",
        "结合历年考点分布和 7 天 deadline...",
    ]
    tf = add_rich_textbox(slide, 1.75, 3.1, 4.8, 1.5)
    for line in diag_text:
        add_para(tf, line, font_size=11, font_color=DARK_TEXT if line.startswith("✓") else GRAY,
                 bold=line.startswith("✓"), space_after=3)

    # Moment 2: Triage Decision (the wow moment)
    add_circle(slide, 0.9, 4.6, 0.3, ORANGE)
    add_rounded_rect(slide, 1.55, 4.3, 5.2, 2.2, SOFT_ORANGE, ORANGE, 2)
    add_textbox(slide, 1.75, 4.4, 4.8, 0.3, "时刻 2：主动放弃",
                font_size=15, font_color=NAVY, bold=True)
    triage_text = [
        "第 8 章考分占比低（~8%），学习成本高（~20h）",
        "投入产出比太低 → 直接放弃",
        "把时间集中在第 4-6 章 + 历年真题",
        "",
        "→ 生成最小通过路径：三块最高转分内容",
    ]
    tf = add_rich_textbox(slide, 1.75, 4.75, 4.8, 1.5)
    for line in triage_text:
        add_para(tf, line, font_size=11, font_color=DARK_TEXT, bold=line.startswith("→"), space_after=3)

    # Right side: visual comparison
    right_x = 7.5
    # "Normal plan" card
    add_rounded_rect(slide, right_x, 2.1, 5.0, 2.0, RGBColor(0xF0, 0xF0, 0xF5))
    add_textbox(slide, right_x + 0.3, 2.2, 4.4, 0.3, "普通 AI 的复习计划",
                font_size=13, font_color=RED, bold=True)
    normal_plan = "第1章 概述  ·  第2章 物理层  ·  第3章 数据链路层\n第4章 网络层  ·  第5章 传输层  ·  第6章 应用层\n第7章 网络安全  ·  第8章 多媒体网络\n\n→ 平均用力，7 天学完 8 章 = 每章不到 1 天"
    add_textbox(slide, right_x + 0.3, 2.6, 4.4, 1.3, normal_plan,
                font_size=10, font_color=GRAY)

    # "Sparkle plan" card
    add_rounded_rect(slide, right_x, 4.4, 5.0, 2.3, SOFT_GREEN, GREEN, 2)
    add_textbox(slide, right_x + 0.3, 4.5, 4.4, 0.3, "Sparkle 的最小通过路径",
                font_size=13, font_color=GREEN, bold=True)
    sparkle_plan = "重点 1  第 4-5 章（网络层 + 传输层） → 35% 考分\n重点 2  第 6 章（应用层） + 历年真题 → 25% 考分\n重点 3  第 1-3 章快速复习 + 错题强化 → 20% 考分\n放弃     第 7-8 章 → 8% 考分，学习成本过高\n\n→ 目标：覆盖 80% 考分，投入产出比最大化"
    add_textbox(slide, right_x + 0.3, 4.9, 4.4, 1.6, sparkle_plan,
                font_size=10, font_color=DARK_TEXT)

    add_speaker_notes(slide, """[50 秒 · 全 deck 最重要的一页之一]

口播：
"我们来看 Sparkle 具体怎么做。

第 0 天。学生做 12 分钟诊断。Sparkle 发现：1 到 3 章还行，4 到 6 章薄弱，7 到 8 章几乎零基础。

接下来是第一个关键决策。Sparkle 分析完考试权重后发现：第 8 章只占约 8% 考分，但学懂需要将近 20 个小时。投入产出比太低。

所以它做了一个普通 AI 不会做的事：主动放弃第 8 章。

不是给学生一份平均复习的全书计划。而是把 7 天时间集中砸在三块最高转分的内容上。

左边是普通 AI 的复习计划：8 章全学，每章不到一天。
右边是 Sparkle：放弃低收益，聚焦 80% 考分。

这不是"Sparkle 更聪明"，而是它有决策模型。它知道什么时候该放弃。"

过渡：但好的决策只是开始。执行阶段出了问题怎么办？""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: RESCUE PART 2 — Execution & Failure Replan
# ═══════════════════════════════════════════════════════════════
def build_slide_5():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "7 天救急 · 下：任务卡 + 失败重规划",
                  "失败不是结束，是系统更新的开始",
                  page_num=5)

    # LEFT: Task card mockup
    left_x = 0.6
    add_textbox(slide, left_x + 0.2, 1.5, 4.0, 0.4, "每天的 Task Card：不是 todo，是执行协议",
                font_size=14, font_color=NAVY, bold=True)

    # Task card visual
    card = add_rounded_rect(slide, left_x, 2.0, 4.5, 3.8, WHITE, LIGHT_GRAY, 2)
    add_shape(slide, left_x, 2.0, 4.5, 0.06, BLUE)

    # Task header
    add_textbox(slide, left_x + 0.25, 2.15, 4.0, 0.35, "Day 3 · 网络层核心协议",
                font_size=13, font_color=NAVY, bold=True)

    sections = [
        ("🎯 为什么做", "IP/子网划分考分占比最高，你的薄弱点"),
        ("📚 用什么", "教材第 4 章 + 2024 真题卷 A"),
        ("📋 怎么做", "① 看 worked example (15min)\n② 做 3 道变式题 (20min)\n③ 小测验证 (10min)"),
        ("🆘 卡住怎么办", "先看解题录像；还不会 → 标记，系统调整下一张卡"),
    ]
    y = 2.6
    for title, body in sections:
        add_textbox(slide, left_x + 0.25, y, 4.0, 0.25, title,
                    font_size=9, font_color=BLUE, bold=True)
        add_textbox(slide, left_x + 0.25, y + 0.22, 4.0, 0.5, body,
                    font_size=8, font_color=DARK_TEXT)
        y += 0.7

    # RIGHT: Failure → Replan flow
    right_x = 5.8
    add_textbox(slide, right_x + 0.2, 1.5, 6.5, 0.4, "第 3 天晚：任务未完成 → 不是打叉，是追问",
                font_size=14, font_color=NAVY, bold=True)

    # Flow: Failure → Attribution → Change
    steps = [
        ("未完成", "25 分钟刷题任务\n未完成", RED, RGBColor(0xFE, 0xE2, 0xE2)),
        ("归因", "系统追问：为什么？\n用户：\"不是没时间\n是不会子网划分\"", ORANGE, SOFT_ORANGE),
        ("改判", "从\"拖延\"改为\n\"知识缺口\"", BLUE, SOFT_BLUE),
        ("重规划", "下一张卡自动换\nworked example\n+ 子网专项练习", GREEN, SOFT_GREEN),
    ]

    for i, (label, desc, accent, bg) in enumerate(steps):
        x = right_x + i * 1.6
        add_rounded_rect(slide, x, 2.1, 1.45, 2.5, bg, accent, 1.5)
        add_textbox(slide, x + 0.1, 2.2, 1.25, 0.25, label,
                    font_size=11, font_color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.6, 1.25, 1.8, desc,
                    font_size=8, font_color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(slide, x + 1.45, 3.0, 0.2, 0.3, "→",
                        font_size=14, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom insight
    add_shape(slide, 0.6, 6.1, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 6.3, 11.5, 0.5,
                "普通工具：任务未完成 = \"失败\"。Sparkle：任务未完成 → 归因 → 改判 → 下一张卡自动变化。",
                font_size=15, font_color=NAVY, bold=True)
    add_textbox(slide, 0.8, 6.75, 11.5, 0.3,
                "必要时，任务承诺可被责任伙伴见证 → 社群共调节（详见附录）",
                font_size=11, font_color=GRAY)

    add_speaker_notes(slide, """[50 秒 · 全 deck 灵魂页]

口播：
"第 3 天，学生拿到一张任务卡——不是 todo，是一份执行协议。它告诉你：为什么做这个任务、用什么资料、分几步做、卡住了怎么办、做完后系统会更新什么。

第 3 天晚上，任务没完成。
在普通 AI 工具里，这就打一个红叉，结束了。
在 Sparkle 里，系统先追问：为什么没做？

学生说：不是没时间，是我不会子网划分。
这就是关键。系统改判——从'拖延'改成'知识缺口'。
然后下一张任务卡自动变：原来是刷题，现在改成 worked example + 子网划分专项练习。

这是一个闭环。失败不是结束，是系统更新的输入。
这也是为什么我们说 Sparkle 不只回答问题，而是在推进目标。"

社群标签一句带过，不展开。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: LONG-TERM VALUE — Personal Method Assets
# ═══════════════════════════════════════════════════════════════
def build_slide_6():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "用一年后，留下的不是聊天记录",
                  "Sparkle 积累的是\"一个人如何做成事\"的证据与方法",
                  page_num=6)

    # Core idea
    add_textbox(slide, 0.8, 1.4, 11.5, 0.5,
                "ChatGPT 用一年：每次对话从零开始。Sparkle 用一年：关于你的成长资产持续累积。",
                font_size=16, font_color=DARK_TEXT)

    # Asset cards
    assets = [
        ("📏", "任务粒度", "适合 25-40 分钟\n的任务卡，不是\n2 小时大任务"),
        ("📚", "资料有效性", "哪些资料真的帮你\n学会了；哪些只是\n看起来有用"),
        ("🔁", "错因模式", "你的失败通常来自\n知识前置缺失，而\n不是拖延或偷懒"),
        ("⚡", "高压策略", "考前先刷真题再补\n理论，比先看教材\n效率高 3 倍"),
        ("🔄", "策略迁移", "计网考试用过的有\n效策略，自动推荐\n到数据库考试"),
        ("🤝", "社群资产", "责任伙伴见证了你\n几次关键突破；群\n体错因帮你避坑"),
    ]

    for i, (icon, title, desc) in enumerate(assets):
        x = 0.4 + i * 2.1
        y = 2.2
        card = add_rounded_rect(slide, x, y, 1.9, 2.8, CARD_BG, LIGHT_GRAY, 1)
        add_textbox(slide, x + 0.15, y + 0.2, 1.6, 0.4, icon,
                    font_size=28, font_color=BLUE)
        add_textbox(slide, x + 0.15, y + 0.7, 1.6, 0.35, title,
                    font_size=12, font_color=NAVY, bold=True)
        add_textbox(slide, x + 0.15, y + 1.2, 1.6, 1.4, desc,
                    font_size=9, font_color=DARK_TEXT)

    # Knowledge star map hint
    add_shape(slide, 0.6, 5.5, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 5.8, 11.5, 0.5,
                "这些不是用户自己总结的——是 Aurora 在每一次决策中自动积累、知识星图持续沉淀的。",
                font_size=14, font_color=NAVY, bold=True)
    add_textbox(slide, 0.8, 6.3, 11.5, 0.4,
                "短期，帮你过眼前这一关。长期，沉淀你怎么变强。",
                font_size=16, font_color=BLUE, bold=True)

    add_speaker_notes(slide, """[45 秒]

口播：
"考试通过。故事结束了吗？没有。系统记住了这一次所有的有效策略。

ChatGPT 用一年，每次对话从零开始。Sparkle 用一年，关于你的成长资产在持续累积。它会知道六件事：

第一，最适合你的任务粒度是 25 到 40 分钟，不是 2 小时大任务。
第二，哪些资料真的帮到你了，哪些只是看起来有用。
第三，你失败通常是因为知识前置缺失，而不是懒。
第四，考前先刷真题再补理论，比先看教材效率高得多。
第五，在计网考试中验证有效的策略，会自动推荐到你的数据库考试。
第六，你的责任伙伴见证了你几次关键突破。

这些不是用户自己总结的。是 Aurora 在每次决策中自动积累的，知识星图持续沉淀的。

一句话：短期，帮你过眼前这一关。长期，沉淀你怎么变强。"

过渡：好，现在你可能在想：为什么这些 ChatGPT 做不到？""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: COMPETITION — Why Others Can't
# ═══════════════════════════════════════════════════════════════
def build_slide_7():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "为什么竞品难做到：不是单点弱，是缺闭环",
                  "现有 AI 优化单次交互质量。Sparkle 优化目标完成过程。",
                  page_num=7)

    # Comparison matrix
    # Header row
    headers = ["", "他们强在哪", "缺什么", "Sparkle 怎么补"]
    col_widths = [3.0, 3.2, 3.2, 3.2]
    header_x = [0.4, 3.5, 6.8, 10.1]

    # Draw header
    for j, (h, x, w) in enumerate(zip(headers, header_x, col_widths)):
        bg_color = BLUE if j > 0 else WHITE
        text_color = WHITE if j > 0 else DARK_TEXT
        add_shape(slide, x, 1.4, w, 0.45, bg_color)
        add_textbox(slide, x + 0.15, 1.43, w - 0.3, 0.4, h,
                    font_size=12, font_color=text_color, bold=True)

    # Data rows
    rows = [
        ("通用大模型\n(ChatGPT/Claude\n/DeepSeek)",
         "最强单次推理\n与内容生成",
         "用户自己管理\n目标与上下文",
         "Aurora 自动组装\n上下文 + 目标闭环",
         2.0),
        ("通用助手\n(豆包等)",
         "易用、入口低\n学习辅助广",
         "偏问答/搜索\n/文档处理",
         "任务执行 + 错因\n反馈 + 策略学习",
         2.6),
        ("解题/学习平台\n(作业帮/夸克等)",
         "知识点和题\n库覆盖强",
         "难形成跨时间\n个体策略",
         "错因 → 重规划\n→ 策略记忆",
         2.3),
        ("校内 AI 平台", "资源和场景\n贴近学生",
         "常是资源/问答\n入口型产品",
         "个体化目标推进\n系统，非入口",
         2.0),
    ]

    y = 2.0
    for row_data, strong, gap, sparkle, h in rows:
        x_positions = [0.4, 3.5, 6.8, 10.1]
        contents = [row_data, strong, gap, sparkle]
        for j, (content, x, w) in enumerate(zip(contents, x_positions, col_widths)):
            bg = CARD_BG if j > 0 else WHITE
            add_shape(slide, x, y, w, h, bg, LIGHT_GRAY, 0.5)
            is_bold = (j == 0)
            fc = NAVY if j == 0 else DARK_TEXT
            add_textbox(slide, x + 0.15, y + 0.1, w - 0.3, h - 0.2, content,
                        font_size=9, font_color=fc, bold=is_bold)
        y += h + 0.08

    # Bottom: Trust dimension
    add_shape(slide, 0.4, y + 0.15, 12.5, 0.015, BLUE)
    add_textbox(slide, 0.6, y + 0.4, 12.0, 0.4,
                "🔒 可信维度：通用 AI 给建议不解释为什么；Sparkle 关键决策可追溯、可纠正、可回滚",
                font_size=14, font_color=NAVY, bold=True)
    add_textbox(slide, 0.6, y + 0.85, 12.0, 0.4,
                "他们优化单点能力。Sparkle 优化目标完成过程。这是两个不同的赛道。",
                font_size=13, font_color=GRAY)

    add_speaker_notes(slide, """[50 秒]

口播：
"为什么现有产品做不到这些？不是因为它们不够好，而是因为它们设计的目标不同。

通用大模型——ChatGPT、Claude、DeepSeek——单次推理和生成能力是天花板级别的。但用户要自己管理目标和上下文。Sparkle 用 Aurora 自动组装上下文，加上目标闭环。

豆包这类通用助手，易用性很强，但更偏向问答搜索。Sparkle 做的是任务执行、错因反馈和策略学习。

解题和学习平台知识点覆盖强，但难形成跨时间的个体策略。

校内 AI 平台资源和场景近，但往往是资源入口，不是个体化目标推进系统。

还有一个维度——可信。通用 AI 给了建议不解释为什么。Sparkle 每次关键决策可追溯、可纠正、可回滚。

一句话：他们优化单点能力。我们优化目标完成过程。这是两个不同的赛道。"

不说"别人做不到"，说"他们优化单点，我们优化过程"。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: ENGINEERING EVIDENCE
# ═══════════════════════════════════════════════════════════════
def build_slide_8():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "不是概念：已有可运行系统",
                  "24 种极端场景全部通过。不是 PPT 项目。",
                  page_num=8)

    # Evidence cards
    evidences = [
        ("24/24", "极端场景全通过", "零基础学生、多目标冲突、\n危机模式、疲劳模式\n—— 全部覆盖并通过", BLUE, SOFT_BLUE),
        ("12/12", "ExamSprint 专项", "7 天考试冲刺全链路\n从诊断到考试结果追踪\n—— 场景全部验证", GREEN, SOFT_GREEN),
        ("97.3%", "完全体达标率", "72/74 验收项 verified\n核心 guard 64/64 全绿\n0 移动端黑屏 bug", ORANGE, SOFT_ORANGE),
        ("3层", "全栈已搭建", "Flutter Mobile App\nGo Gateway 统一调度\nPython AI Engine + gRPC", RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xED, 0xE9, 0xFE)),
    ]

    for i, (number, title, desc, accent, bg) in enumerate(evidences):
        x = 0.4 + i * 3.15
        y = 1.6
        add_rounded_rect(slide, x, y, 2.95, 2.8, bg, accent, 1.5)
        add_textbox(slide, x + 0.2, y + 0.2, 2.55, 0.7, number,
                    font_size=42, font_color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, y + 0.95, 2.55, 0.35, title,
                    font_size=16, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, y + 1.4, 2.55, 1.2, desc,
                    font_size=10, font_color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Product screenshot placeholder
    add_textbox(slide, 0.8, 4.8, 5.0, 0.35, "📱 产品实机截图",
                font_size=14, font_color=NAVY, bold=True)

    # Phone frame placeholders
    for i in range(3):
        x = 0.8 + i * 2.0
        phone = add_phone_frame(slide, x, 5.2, 1.8, 3.6)
        # Screen placeholder
        add_rounded_rect(slide, x + 0.15, 5.35, 1.5, 2.3, LIGHT_BG)
        labels = ["任务卡界面", "知识星图", "对话决策"]
        add_textbox(slide, x + 0.15, 6.5, 1.5, 0.3, labels[i],
                    font_size=8, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Note
    add_textbox(slide, 7.0, 5.0, 5.8, 0.3,
                "▸ 后续可替换为真实产品截图",
                font_size=9, font_color=GRAY)

    add_speaker_notes(slide, """[40 秒]

口播：
"到这里你可能在想：这真的做出来了吗？
是的。Sparkle 不是 PPT 项目。

我们用 24 种极端场景测试了系统——包括零基础学生、多目标冲突、危机模式、疲劳模式。全部通过。
考试冲刺全链路——从诊断到考试结果追踪——12 个场景全部验证。
72 个验收项里 74 个 checkpoint，97.3% 通过。核心 guard 64 个全绿。移动端零黑屏 bug。
三端全栈已搭建：Flutter 移动端、Go 网关、Python AI 引擎。

这是我们产品的真实截图——任务卡、知识星图、对话决策界面。后续会替换为最新版本。"

暗示力量：我们是能交付的团队。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: TEAM & ADVISOR
# ═══════════════════════════════════════════════════════════════
def build_slide_9():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "为什么我们能做成：团队与导师",
                  "这个项目需要的不只是编程能力，而是五层能力栈的同时具备",
                  page_num=9)

    # Advisor
    add_textbox(slide, 0.8, 1.5, 5.0, 0.35, "项目导师",
                font_size=14, font_color=BLUE, bold=True)
    add_rounded_rect(slide, 0.8, 1.95, 5.0, 1.6, SOFT_BLUE, BLUE, 1.5)
    add_textbox(slide, 1.1, 2.1, 4.4, 0.35, "黄亚坤  副教授 · 博士",
                font_size=16, font_color=NAVY, bold=True)
    add_textbox(slide, 1.1, 2.5, 4.4, 0.3, "北京邮电大学计算机学院（国家示范性软件学院）",
                font_size=11, font_color=DARK_TEXT)
    add_textbox(slide, 1.1, 2.85, 4.4, 0.5, "研究方案指导 · 技术路线把关\n实验设计 · 论文写作指导",
                font_size=11, font_color=GRAY)

    # Team capability stack
    add_textbox(slide, 6.5, 1.5, 5.0, 0.35, "团队能力栈",
                font_size=14, font_color=BLUE, bold=True)

    capabilities = [
        ("💡", "产品洞察", "为什么选 7 天考试救急\n作为北极星压力测试"),
        ("🧠", "AI 系统", "Aurora、策略闭环\n用户模型、因果控制链"),
        ("📱", "全栈工程", "Flutter + Go Gateway\n+ Python AI Engine"),
        ("🎓", "教育场景", "考试冲刺、学习科学\n可验证性把关"),
        ("🛡️", "安全治理", "53+ 治理规则\n全链路可追溯"),
    ]

    for i, (icon, title, desc) in enumerate(capabilities):
        x = 6.5 + (i % 3) * 2.1
        y = 2.0 + (i // 3) * 1.8
        add_rounded_rect(slide, x, y, 1.9, 1.6, CARD_BG, LIGHT_GRAY, 1)
        add_textbox(slide, x + 0.15, y + 0.15, 1.6, 0.3, icon + " " + title,
                    font_size=12, font_color=NAVY, bold=True)
        add_textbox(slide, x + 0.15, y + 0.55, 1.6, 0.9, desc,
                    font_size=9, font_color=DARK_TEXT)

    # Team members
    add_textbox(slide, 0.8, 3.9, 5.0, 0.35, "团队成员",
                font_size=14, font_color=BLUE, bold=True)

    members = [
        ("邓博仁", "2006 级 · 计算机科学与技术", "产品愿景 / 系统架构"),
        ("张雨凝", "2005 级 · 计算机科学与技术", "AI 系统 / 后端工程"),
        ("王  宇", "2006 级 · 电子信息", "移动端 / 体验设计"),
        ("王英树", "2006 级 · 计算机科学与技术", "数据闭环 / 测试验证"),
    ]

    for i, (name, info, role) in enumerate(members):
        x = 0.8 + i * 3.1
        add_rounded_rect(slide, x, 4.35, 2.9, 1.3, CARD_BG, LIGHT_GRAY, 1)
        add_textbox(slide, x + 0.2, 4.45, 2.5, 0.3, name,
                    font_size=15, font_color=NAVY, bold=True)
        add_textbox(slide, x + 0.2, 4.75, 2.5, 0.25, info,
                    font_size=9, font_color=GRAY)
        add_textbox(slide, x + 0.2, 5.05, 2.5, 0.4, role,
                    font_size=10, font_color=BLUE)

    # Bottom line
    add_shape(slide, 0.6, 6.0, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 6.2, 11.5, 0.5,
                "这个项目不是一个人靠 prompt 做出来的。它需要产品洞察、AI 系统设计、全栈工程、教育场景理解和安全治理五层能力的交叉。",
                font_size=14, font_color=DARK_TEXT)

    add_speaker_notes(slide, """[40 秒]

口播：
"为什么我们能把这个很难的系统做出来？
因为这个项目需要的不是单一能力，而是五层能力的交叉。

产品洞察：为什么选 7 天考试救急作为北极星，而不是做一个通用的学习问答？
AI 系统：Aurora、策略闭环、用户模型、因果控制链——这些不是 prompt engineering。
全栈工程：Flutter 移动端、Go 网关、Python AI 引擎，三端全栈。
教育场景：考试冲刺是真实刚需，需要学习科学的判断。
安全治理：53+ 条治理规则，全链路可追溯。

我们的导师黄亚坤副教授，北邮计算机学院，在研究方法、技术路线和实验设计上全程把关。

团队四个人，各有分工。这不是一个人靠 prompt 拼出来的项目。"

不说"我们很强"，说"这个项目需要这五种能力，我们团队恰好覆盖"。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: FUTURE — From Exam Rescue to Goal OS
# ═══════════════════════════════════════════════════════════════
def build_slide_10():
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "从校园考试救急，到 AI-native Goal OS",
                  "先在一个刚需场景做到极致，再系统性扩展",
                  page_num=10)

    # Expansion path
    phases = [
        ("现在", "北邮试点\n考试冲刺 PMF", BLUE, SOFT_BLUE),
        ("近期", "更多学科 · 更多高校\nC端订阅 + 机构合作", GREEN, SOFT_GREEN),
        ("中期", "项目交付 · 求职\n科研 · 跨目标迁移", ORANGE, SOFT_ORANGE),
        ("长期", "AI-native Goal OS\n任何目标都能推进", RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xED, 0xE9, 0xFE)),
    ]

    for i, (label, desc, accent, bg) in enumerate(phases):
        x = 0.5 + i * 3.2
        add_rounded_rect(slide, x, 1.5, 2.95, 2.2, bg, accent, 1.5)
        add_textbox(slide, x + 0.2, 1.65, 2.55, 0.35, label,
                    font_size=16, font_color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, 2.15, 2.55, 1.2, desc,
                    font_size=12, font_color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(slide, x + 2.95, 2.2, 0.3, 0.4, "→",
                        font_size=22, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # Business model - brief
    add_shape(slide, 0.6, 4.2, 12.1, 0.015, BLUE)
    add_textbox(slide, 0.8, 4.4, 11.5, 0.35, "商业模式",
                font_size=14, font_color=BLUE, bold=True)

    # Two columns
    add_rounded_rect(slide, 0.8, 4.85, 5.5, 1.5, SOFT_BLUE)
    add_textbox(slide, 1.1, 5.0, 5.0, 0.3, "C 端订阅",
                font_size=14, font_color=NAVY, bold=True)
    add_textbox(slide, 1.1, 5.35, 5.0, 0.8,
                "个人用户付费会员\n增值功能（深度诊断、策略报告、跨目标迁移）\n考试是刚需 · 学生愿意为通过付费",
                font_size=10, font_color=DARK_TEXT)

    add_rounded_rect(slide, 7.0, 4.85, 5.5, 1.5, SOFT_GREEN)
    add_textbox(slide, 7.3, 5.0, 5.0, 0.3, "D 端合作（B2B）",
                font_size=14, font_color=NAVY, bold=True)
    add_textbox(slide, 7.3, 5.35, 5.0, 0.8,
                "与教育机构、高校开展深度合作\n企业培训场景扩展\n盈利优先 · 健康增长 · 拒绝烧钱",
                font_size=10, font_color=DARK_TEXT)

    # Policy alignment
    add_textbox(slide, 0.8, 6.6, 11.5, 0.35,
                "深度契合国家\"人工智能+\"行动与教育数字化战略",
                font_size=13, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide, """[35 秒]

口播：
"我们的路径很清晰。现在：北邮试点，用考试冲刺验证 PMF。
近期：扩展到更多学科、更多高校。C 端订阅加机构合作。
中期：从考试延伸到项目交付、求职、科研——跨目标迁移。
长期：AI-native Goal OS——不只是学习，是任何目标都能推进。

商业上两条线：C 端订阅，考试是刚需，学生愿意为通过付费。D 端机构合作，企业培训场景。

一句话：深度契合国家人工智能+行动和教育数字化战略。"

快过，不展开。留给收束。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: CLOSE — Back to the Student
# ═══════════════════════════════════════════════════════════════
def build_slide_11():
    slide = add_blank_slide()
    set_slide_bg(slide, NAVY)

    # Return to the student
    add_textbox(slide, 0.8, 1.0, 11.5, 0.6, "7 天前",
                font_size=18, font_color=GRAY)
    add_textbox(slide, 0.8, 1.5, 11.5, 0.8,
                "那个面对计网考试、基本没学的学生",
                font_size=28, font_color=WHITE)

    add_textbox(slide, 0.8, 2.5, 11.5, 0.6, "7 天后",
                font_size=18, font_color=GRAY)
    add_textbox(slide, 0.8, 3.0, 11.5, 0.8,
                "他通过了。",
                font_size=36, font_color=ORANGE, bold=True)

    add_textbox(slide, 0.8, 3.6, 11.5, 0.5,
                "不是因为 Sparkle 比他聪明。",
                font_size=18, font_color=GRAY)

    add_textbox(slide, 0.8, 4.1, 11.5, 0.5,
                "而是因为 Sparkle 帮他做了那些他不知道该怎么做的决策。",
                font_size=22, font_color=WHITE, bold=True)

    # Three takeaways
    takeaway_x = 0.8
    takeaways = [
        "Sparkle 不只是答题，而是帮学生把目标做成。",
        "7 天考试救急，是我们证明价值的第一战场。",
        "长期来看，它会沉淀每个人自己的成长方法。",
    ]
    for i, t in enumerate(takeaways):
        add_textbox(slide, takeaway_x, 5.1 + i * 0.45, 11.5, 0.4,
                    f"{i+1}.  {t}", font_size=16, font_color=WHITE)

    # Final logo line
    add_shape(slide, 0.6, 6.6, 2.0, 0.02, ORANGE)
    add_textbox(slide, 0.8, 6.8, 8.0, 0.5,
                "Sparkle  ·  让 AI 不只会讲题，而是真的陪你把目标走完",
                font_size=20, font_color=WHITE, bold=True)

    add_speaker_notes(slide, """[30 秒 · 情绪收束]

口播：
"让我们回到那个学生。7 天前，他面对计网考试，基本没学。7 天后，他通过了。
不是因为 Sparkle 比他聪明。
而是因为 Sparkle 帮他做了那些他不知道该怎么做的决策——该救哪三块、该放弃什么、资料怎么用、失败后下一步变什么。

三句话，我们希望评委记住：
第一，Sparkle 不只是答题，而是帮学生把目标做成。
第二，7 天考试救急，是我们证明价值的第一战场。
第三，长期来看，它会沉淀每个人自己的成长方法。

Sparkle——让 AI 不只会讲题，而是真的陪你把目标走完。

谢谢。"

最后 10 秒留白。让评委记住那个从恐慌到通过的学生。""")
    return slide


# ═══════════════════════════════════════════════════════════════
# APPENDIX SLIDES (8 pages for Q&A defense)
# ═══════════════════════════════════════════════════════════════

def build_appendix(title, bullets, notes=""):
    slide = add_blank_slide()
    set_slide_bg(slide, WHITE)
    # Appendix label
    add_shape(slide, 0, 0, 13.333, 0.06, ORANGE)
    add_textbox(slide, 0.8, 0.3, 3.0, 0.35, "附 录",
                font_size=10, font_color=ORANGE, bold=True)
    add_textbox(slide, 0.8, 0.8, 11.5, 0.6, title,
                font_size=28, font_color=NAVY, bold=True)

    tf = add_rich_textbox(slide, 0.8, 1.6, 11.5, 5.0)
    for bullet in bullets:
        add_para(tf, bullet, font_size=14, font_color=DARK_TEXT, space_after=10)

    if notes:
        add_speaker_notes(slide, notes)
    return slide


def build_all_appendix():
    slides_data = [
        ("A. 社群责任伙伴机制", [
            "• 责任伙伴闭环：承诺 → 伙伴提醒 → 见证 → 外部观察候选 → 用户确认",
            "• 群聊与火堆（社群空间）：基于共同目标的临时社群，非社交网络",
            "• 社群共性错因：群体级别的错因聚类 → 推荐群体级策略调整",
            "• 伙伴信号写入 ActionableStatePacket（需用户确认，遵守隐私铁律）",
            "• 社群资源质量评分：下载量不是标准，outcome + 负反馈 + 适用范围才是",
            "• 铁律约束：伙伴观察仅作 external_observation_candidate + needs_user_confirmation",
        ], "答辩问题：'你们只有 AI 对话吗？社群是怎么做的？'"),
        ("B. 安全、隐私、可追溯体系", [
            "• 身份与访问控制：JWT + WebSocket 票据 + 白名单 + RBAC",
            "• 全链路 Request ID 追踪：每一步决策可追溯到输入和规则",
            "• 用户纠正直达底层：用户矫正绕过 Aurora，写入 L0 raw evidence",
            "• 反证降权（Counter-Evidence）：3 次反证后 belief 自动失效",
            "• Aurora 白名单纪律：调旋钮不改电路（数据源、算法、证据链不可被 Aurora 修改）",
            "• 53+ 治理规则 CI 强制：每次 PR 合并前自动运行",
            "• Kill Switch 协议：所有 Aurora 特性 tri-state（off → shadow → live），异常自动降级",
            "• 数据最小化 fail-closed：15 个 scope，GOV-DATA-MIN guard 强制执行",
        ], "答辩问题：'AI 的建议可信吗？出了问题谁负责？数据安全怎么保证？'"),
        ("C. Aurora 分层与交互式建模", [
            "• Aurora 技术本质：Adaptive Context Engineering — 每次交互动态组装真正匹配当下的 context",
            "• 三时层模型（inline / nearline / long-horizon）：同一引擎在三种预算约束下的行为谱",
            "  - inline: P95 < 100ms，同步路径，BackboneRoutingDecision",
            "  - nearline: P95 ≤ 30s，会话结束/idle/提交，TDR/InsightClaim/ProbeOutcome",
            "  - long-horizon: 小时~天级，批处理，FocusContract 演进/IdentityEvidence 聚合",
            "• 交互式校准：用户通过对话让系统理解自己，系统保持独立判断，不完全听信用户",
            "• 双交互模式：普通对话（低 Aurora 介入） vs 深度模式（高 Aurora 介入，走进内心世界）",
            "• 五层用户模型（L0-L3 + 用户纠偏）：层层写入隔离，用户纠偏直达 raw/calibration",
            "• Aurora 影子模型：轻量状态摘要，漂移检测——Aurora 说的和画像系统说的不一致 = 需要问用户",
        ], "答辩问题：'Aurora 到底是什么？它不是又一个推荐算法吗？'"),
        ("D. 知识星图与个人知识库", [
            "• 知识星图不是展示图，而是目标世界模型：资料、考点、错因、任务、考试权重都挂在上面",
            "• 资料不是附件，而是世界模型的一部分：上传后挂到星图节点、任务卡、错因和计划里",
            "• 6+1 知识星图体系：课程知识 / 考试考点 / 错因模式 / 学习策略 / 资料有效性 / 个人方法 + 社群知识",
            "• 星图溯源：每次决策可追溯到用了哪些资料节点、为什么选这些、效果如何",
            "• Source Tray + Context Receipt：资料使用有 trace——用了什么、没用什么、为什么",
            "• 长期资产：星图随使用持续生长，跨目标可查询、可复用",
            "• 技术栈：PostgreSQL + pgvector + Apache AGE（图数据库）+ Redis 缓存",
        ], "答辩问题：'个人知识管理怎么做的？和 Notion/文档管理有什么区别？'"),
        ("E. Skill Extraction / Learning Base", [
            "• Skill 三级体系：personal（个人有效策略）→ cohort（群体验证）→ system（平台级资产）",
            "• 自动提取触发：不只 outcome_positive，多条件触发（采纳→完成→反馈正→复用→再正）",
            "• Learning Base：贝叶斯信念更新 + 冷启动先验 + 后验更新 + 长期信念保留（跨 sprint）",
            "• Skill 生命周期：inject → extract → recommend → verify → promote/deprecate",
            "• 晋升门槛：personal→cohort 需 5 次有效验证；cohort→system 需 10 次多用户验证",
            "• 反证降权：counter_evidence 累积 3 次 → belief 失效 → skill 降级/下架",
            "• Marketplace 上线：Skill 可预览、采纳、追踪效果、负反馈自动下架",
            "• 长期价值：Sparkle 越用越懂你，不是比喻——是可验证的策略积累和自动推荐",
        ], "答辩问题：'系统怎么越用越好？你说是学习，具体怎么学的？'"),
        ("F. MirrorFish / 多 Agent 推演", [
            "• MirrorFish 推演系统：角色扮演 + 多 Agent 活动，模拟真实目标推进场景",
            "• 多 Agent 场景：学生 Agent + 导师 Agent + 伙伴 Agent + 压力 Agent（模拟考试焦虑）",
            "• 策略预验证：新策略先在 MirrorFish 推演，通过后再进入 shadow → canary → safe_live",
            "• Simulated Gray Window (SGW)：预发布前 12h + 200 sessions 模拟验证",
            "• 24 场景基准测试（SparkleGoalBench）：ExamSprint 12 + ProjectDelivery 4 + JobSearch 4 + MultiGoalLife 4",
            "• SyntheticPersona 模拟 + TraceReplay 回放：混合策略，不作为唯一评估依据",
            "• 极端场景必过：D0 危机 / 零基础 / 多目标冲突 / 疲劳模式",
            "• 与 Causal Spine 配合：推演结果写入 outcome → 反事实评估 → 策略更新建议",
        ], "答辩问题：'你们的策略怎么验证有效？怎么保证不会给学生错误的建议？'"),
        ("G. 技术架构全图", [
            "• Flutter (Presentation) → Go Gateway (Coordination) → Python AI Engine (Intelligence)",
            "• Go Gateway：Auth（JWT + WebSocket）、路由、缓存（Redis）、流式推送、限流",
            "• Python AI Engine：LangGraph FSM（Orchestrator）、Dual-Core Router、Aurora Adaptive Kernel",
            "• 数据层：PostgreSQL 16 + pgvector（向量相似）+ Apache AGE（图查询）+ Redis Stack（缓存/事件总线）",
            "• gRPC 内部通信：Go ↔ Python，protobuf 定义接口（6 个 proto 文件）",
            "• WebSocket 实时通道：Flutter ↔ Go ↔ Python，server-streaming",
            "• Celery 异步任务：定时 beat schedule（策略学习、技能提取、反事实评估、guardrail 检查）",
            "• 外部集成：OpenClaw 执行闭环 + 多模型路由（DeepSeek V4 主力 + 10+ 模型自适应调度）",
            "• 代码规模：86 万+行业务代码，13 万+行测试，7342 个测试用例收集",
        ], "答辩问题：'你们的技术方案是什么样的？整体架构介绍一下。'"),
        ("H. 工程验证与测试指标详情", [
            "• SparkleGoalBench：24/24 场景全部通过（ExamSprint 12 + ProjectDelivery 4 + JobSearch 4 + MultiGoalLife 4）",
            "• ExamSprintBench：12/12 专项通过（7 天全链路 + 退化模式 + 社群 + 疲劳 + 危机模式）",
            "• 完全体达标率：72/74 checkpoint verified = 97.3%",
            "• Rule Guards：64/64 核心规则全部通过（含 K/Y/Z/AB/AF/AM-AN-AO-AP-AQ/AS-AT-AU-AV/AW-AX-AY-AZ/BB/BC）",
            "• 后端单测：153 个 FV 目标测试通过，0 失败",
            "• Go Gateway：5 个包全部通过（cqrs / db / handler / middleware / service）",
            "• Alembic 迁移：单一 head (c24_20260502)，无多头分叉",
            "• 总测试收集：7342 tests collected，零 collection error",
            "• 一票否决项：0 / 10 触发",
            "• 本地真机联调：90% 关键路径已跑通",
            "• 北极星指标在线：7 天目标完成率 / 考试通过概率 / 实际考试结果追踪",
            "• 移动端：0 黑屏 bug，Flutter 698 pass / 3 pre-existing fails (IsarCore)",
        ], "答辩问题：'你们做到什么程度了？有什么具体数据吗？'"),
    ]

    for title, bullets, notes in slides_data:
        build_appendix(title, bullets, notes)


# ═══════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════════

build_slide_1()
build_slide_2()
build_slide_3()
build_slide_4()
build_slide_5()
build_slide_6()
build_slide_7()
build_slide_8()
build_slide_9()
build_slide_10()
build_slide_11()
build_all_appendix()

# ── Save ─────────────────────────────────────────────────────
output_path = "/Users/brsama/code/GitHub/Sparkle-project/docs/presentations/Sparkle_V2_Roadshow.pptx"
prs.save(output_path)
print(f"✓ PPT saved to: {output_path}")
print(f"  Slides: {len(prs.slides)} (11 main + 8 appendix)")
